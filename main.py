from flask import Flask, render_template, request
import pandas as pd
import numpy as np
from werkzeug.utils import secure_filename
import cv2
from pptx import Presentation
from pptx.util import Inches
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
import pythoncom
import win32com.client
import os
from flask import Flask
from flask_pymongo import PyMongo
import time
import qrcode
import random


def make_certificate(filename, template, date, title):
    i = 0
    cid = int(time.time()*random.randint(0, 1000))
    data = pd.read_csv(filename)
    names = np.array(data['Name'])
    for name in names:
        prs = Presentation(template)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if (shape.text.find("{{FULL_NAME}}")) != -1:
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                cur_text = run.text
                                i += 1
                                new_text = cur_text.replace(
                                    "{{FULL_NAME}}", str(name))
                                run.text = new_text
                    if (shape.text.find("{{DATE}}")) != -1:
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                cur_text = run.text
                                i += 1
                                new_text = cur_text.replace(
                                    "{{DATE}}", str(date))
                                run.text = new_text
                    # if (not mongo.db.certificates.find_one({'name': name}) and mongo.db.certificates.find_one({'Occasion': title})):
                        certificate = {'name': name,
                                       'Occasion': title,
                                       'date': date,
                                       'cid': cid}
                        mongo.db.certificates.insert_one(certificate)
                    qr = qrcode.QRCode(version=None, box_size=1)
                    data = f"https://localhost:3000/api/verify/{cid}"
                    qr.add_data(data)
                    qr.make(fit=True)
                    img = qr.make_image(fill_color='black',
                                        back_color='white')
                    img.save(f"static/{name}qrcode.png")
                    img_path = f"static/{name}qrcode.png"
                    pic = slide.shapes.add_picture(
                        img_path, 0, 0)
                    prs.save(f"static/{name}.pptx")
                    i += 1


def convert():
    Path = os.listdir(os.getcwd()+'\\static')
    folder = 'static'
    for file in Path:
        if file.endswith('.pptx'):
            output_file = os.path.splitext(file)[0] + '.pdf'

            # Use pywin32 to convert the PowerPoint file to a PDF
            pythoncom.CoInitialize()
            ppt = win32com.client.Dispatch('PowerPoint.Application')
            presentation = ppt.Presentations.Open(
                os.path.abspath(os.path.join(folder, file)))
            presentation.SaveAs(os.path.abspath(
                os.path.join(folder, output_file)), 32)
            presentation.Close()
            ppt.Quit()
            pythoncom.CoUninitialize()
            os.remove(os.path.join(folder, file))


def send_mail(filename, title):
    data = pd.read_csv(filename)
    names = np.array(data['Name'])
    emails = np.array(data['Email'])

    for i in range(len(emails)):
        fromaddr = "tony@tonyflorida.me"
        toaddr = emails[i]

        msg = MIMEMultipart()

        msg['From'] = fromaddr

        msg['To'] = toaddr

        msg['Subject'] = f"Certificate for {title}"

        body = f"Hello {names[i]}, Here is your certificate for {title}"

        msg.attach(MIMEText(body, 'plain'))

        filename = f"static/{names[i]}.pdf"
        attachment = open(filename, "rb")

        p = MIMEBase('application', 'octet-stream')

        p.set_payload((attachment).read())

        encoders.encode_base64(p)

        p.add_header('Content-Disposition',
                     "attachment; filename= %s" % filename)

        msg.attach(p)

        s = smtplib.SMTP('smtp.dreamhost.com', 587)

        s.starttls()

        s.login('tony@tonyflorida.me', "mystrongpw!")

        text = msg.as_string()

        s.sendmail(fromaddr, toaddr, text)

        s.quit()


app = Flask(__name__)
app.config["MONGO_URI"] = "mongodb://localhost:27017/certificategenerator"
mongo = PyMongo(app)


@app.route("/")
def home():
    return render_template("index.html")


@app.route('/api/<cid>')
def hello(cid):
    certificate = mongo.db.certificates.find_one({'cid': cid})
    if (certificate):
        return render_template('verify.html')
    else:
        return render_template("not-verify.html")


@app.route("/api/file", methods=["GET", "POST"])
def operation():
    if request.method == "POST":
        date = request.form.get('date')
        title = request.form.get('title')
        f1 = request.files['file']
        f2 = request.files['template']
        f1.save(f'uploads/{secure_filename(f1.filename)}')
        f2.save(f'uploads/{secure_filename(f2.filename)}')
        make_certificate(f1.filename, f2.filename, date, title)
        convert()
        send_mail(f1.filename, title)
        return render_template("results.html")


app.run(debug=True)
